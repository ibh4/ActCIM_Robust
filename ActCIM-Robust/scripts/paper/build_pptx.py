"""Build the 12-slide editable PPTX for ActCIM-Robust from the final paper.

Layout per content slide: title bar + one-sentence core conclusion (accent
box) + 3-5 short bullets (left) + the corresponding 300-DPI figure (right).
All numbers are the audited values used in the final paper.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from PIL import Image

ROOT = Path(__file__).resolve().parents[2]
FIG = ROOT / "results" / "figures" / "paper_final"
OUT = ROOT / "reports" / "final" / "ActCIM_Robust_slides.pptx"

EA = "Songti SC"
LAT = "Times New Roman"
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xD6, 0x27, 0x28)
GREY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _set_font(run, size, bold=False, color=GREY, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = LAT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", EA)


def add_textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, shrink=False):
    """lines: list of (text, size, bold, color, bullet_indent_level_or_None).
    shrink=True enables PowerPoint's shrink-text-on-overflow autofit."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if shrink:
        bodyPr = tf._txBody.bodyPr
        for tag in ("a:noAutofit", "a:normAutofit", "a:spAutoFit"):
            el = bodyPr.find(qn(tag))
            if el is not None:
                bodyPr.remove(el)
        bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {}))
    first = True
    for text, size, bold, color, bullet in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(6)
        if bullet is not None:
            p.level = bullet
            run = p.add_run()
            run.text = "▪ " + text
        else:
            run = p.add_run()
            run.text = text
        _set_font(run, size, bold=bold, color=color)
    return tb


def add_rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


def add_picture_fit(slide, img_path: Path, x, y, max_w, max_h):
    """Insert the image centered inside the (x, y, max_w, max_h) inch region."""
    with Image.open(img_path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = iw * ratio, ih * ratio
    slide.shapes.add_picture(str(img_path),
                             Inches(x + (max_w - w) / 2),
                             Inches(y + (max_h - h) / 2),
                             width=Emu(int(w * 914400)),
                             height=Emu(int(h * 914400)))


def content_slide(num, title, conclusion, bullets, fig=None, note=None,
                  fig_region=(6.55, 1.62, 6.45, 5.35), bullet_w=6.15):
    s = prs.slides.add_slide(BLANK)
    # title bar
    add_rect(s, 0, 0, 13.333, 0.92, NAVY)
    add_textbox(s, 0.45, 0.08, 11.6, 0.78,
                [(f"{num}  {title}", 24, True, RGBColor(0xFF, 0xFF, 0xFF), None)],
                anchor=MSO_ANCHOR.MIDDLE)
    # conclusion band
    add_rect(s, 0.45, 1.05, 12.45, 0.52, LIGHT)
    add_textbox(s, 0.62, 1.07, 12.2, 0.5,
                [("核心结论：" + conclusion, 14.5, True, ACCENT, None)],
                anchor=MSO_ANCHOR.MIDDLE)
    # bullets (1.5x size, shrink-on-overflow keeps long items inside)
    lines = [(b, 21, False, GREY, 0) for b in bullets]
    add_textbox(s, 0.45, 1.85, bullet_w, 5.35, lines, shrink=True)
    # figure
    if fig is not None:
        fx, fy, fw, fh = fig_region
        add_picture_fit(s, FIG / fig, fx, fy, fw, fh)
    if note:
        add_textbox(s, 0.45, 6.95, 12.5, 0.45,
                    [(note, 10, False, RGBColor(0x88, 0x88, 0x88), None)])
    return s


# ---------------------------------------------------------------- slide 1
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, NAVY)
add_textbox(s, 1.0, 2.15, 11.3, 1.9, [
    ("面向存算一体激活非线性的鲁棒性建模与非线性感知训练", 30, True,
     RGBColor(0xFF, 0xFF, 0xFF), None),
    ("ActCIM-Robust：CIFAR-10 / ResNet-18 上的系统研究", 20, False,
     RGBColor(0xCF, 0xDD, 0xEE), None),
])
add_textbox(s, 1.0, 4.35, 11.3, 1.6, [
    ("Fixed-NAT(+0.4)：最差准确率 81.25% → 91.79%，AURC 0.9283 → 0.9374，"
     "Clean 保持 94.02%", 16, True, RGBColor(0xFF, 0xD9, 0x66), None),
    ("全部数值可由项目 CSV / 日志 / checkpoint 追溯（附结果审计说明）", 13,
     False, RGBColor(0xB8, 0xC7, 0xDA), None),
])

# ---------------------------------------------------------------- slide 2
content_slide(
    "02", "研究背景：存算一体的激活非线性",
    "CIM 模拟链路给激活引入确定性非线性失真，精度退化不可控且方向未知。",
    [
        "存算一体将乘加下沉到存储阵列，能耗数量级降低（ISAAC/PRIME 等）",
        "驱动放大、位线传输、ADC 有限线性区 → 激活经历非线性传输",
        "既有工作聚焦权重侧噪声；确定性激活非线性关注不足",
        "失真沿网络深度逐层累积，同时破坏精度与置信度校准",
        "目标：建立可追溯的评估协议 + 对比三种非线性感知训练（NAT）",
    ],
    fig="fig01_pipeline.png",
    fig_region=(6.35, 1.72, 6.7, 5.2),
    bullet_w=5.9)

# ---------------------------------------------------------------- slide 3
content_slide(
    "03", "非线性数学模型 f_α",
    "单参数三次族 f_α 统一刻画压缩（α>0）与扩张（α<0），端点不动、处处可导。",
    [
        "f_α(x) = m·[α(x/m)³ + (1−α)(x/m)]，m = max|x| 逐张量动态归一化",
        "α>0：小信号增益 1−α < 1，激活压缩、趋向三次饱和",
        "α<0：小信号增益 > 1，激活扩张但保序性更好",
        "一阶齐次：对激活整体尺度不敏感；可解析反传，无需 STE",
        "注入位置：层输入端 y = W·f_α(x)+b（实际生效深层 4 层，见 P05）",
    ],
    fig="fig02_nonlinearity.png",
    fig_region=(6.3, 1.72, 6.75, 5.2),
    bullet_w=5.85)

# ---------------------------------------------------------------- slide 4
content_slide(
    "04", "技术路线与评估协议",
    "基线 → 非线性注入 → 三种 NAT 微调 → 统一 11 点 α-Sweep 可追溯评估。",
    [
        "ResNet-18-CIFAR（11.18M 参数）基线：50 epoch，val 94.84%",
        "三种 NAT：Fixed(α=+0.4 恒定) / Random(α~U(−0.5,+0.5)) / SGR(灵敏度引导+KL)",
        "评估：α∈[−0.8,+0.8] 11 点 × 10 000 张全测试集，batch 256",
        "指标：Worst-Case Acc、AURC(梯形积分归一化)、Asymmetry Gap、ECE",
        "全部 checkpoint 复评核验：α=0 指纹与 CSV 逐位一致",
    ],
    fig="fig01_pipeline.png",
    fig_region=(6.35, 1.72, 6.7, 5.2),
    bullet_w=5.9)

# ---------------------------------------------------------------- slide 5
content_slide(
    "05", "实验设置与结果审计",
    "审计实证锁定：头条数字全部来自 Fixed-NAT seed 42 checkpoint；注入实际作用于深层 4 层。",
    [
        "seed 口径不一致已解决：val acc 最高为 seed 2026（94.98%），但 α-Sweep "
        "CSV 指纹（α=0 → 94.02%）唯一匹配 seed 42 → 论文统一以 seed 42 为主结果",
        "控制器命名碰撞：21 层登记名塌缩为 4 键，enable_all() 实际启用 "
        "layer4.1.conv1/conv2、layer4.0.downsample.0、fc",
        "扰动强度约为设计意图六成（logit 相对 L2：0.59 vs 0.98）",
        "训练与评估同口径 → 四方法对比公平；结论限定于\u201c深层非线性扰动模型\u201d",
        "除 Fixed-NAT（n=3）外各方法单种子：不宣称多种子统计显著性",
    ],
    fig="fig10_multiseed.png",
    fig_region=(7.3, 1.85, 5.75, 4.95),
    bullet_w=6.75)

# ---------------------------------------------------------------- slide 6
content_slide(
    "06", "方向不对称：压缩远比扩张危险",
    "正 α（压缩）是主导失效方向：基线 α=+0.8 跌至 81.25%，α=−0.8 仍有 93.66%。",
    [
        "基线方向不对称差 −3.14 pp（正侧均值低于负侧）",
        "机理：压缩把中小幅值激活增益压至 1−α，判别信息被系统性抹除",
        "负 α 的激活扰动幅度更大（fc 处相对 L2 0.98 vs 0.61）但保序性好",
        "决定精度损失的不是扰动幅度，而是压缩对判别信息的破坏",
        "Fixed-NAT 训练后不对称差归零（+0.0007）→ 敏感方向可被训练消除",
    ],
    fig="fig05_error_accumulation.png",
    fig_region=(6.55, 1.72, 6.5, 5.2),
    bullet_w=6.05)

# ---------------------------------------------------------------- slide 7
content_slide(
    "07", "层敏感性分析：证据与局限",
    "单层排序实验因命名碰撞与 128 样本分辨率而退化，有效证据来自误差逐层累积。",
    [
        "原设计：逐层单独注入 α=±0.4 对 21 层排序",
        "退化 1：命名碰撞 → 21 条记录仅映射到 4 个唯一生效层，组内数值重复",
        "退化 2：单个 128 样本批，分辨率 0.78 pp，实测最大变化恰为 ±0.78 pp",
        "本文不将其作为逐层敏感性证据；SGR-NAT 的灵敏度先验因此偏弱",
        "有效证据：误差自 layer4.0.downsample.0 起沿深度逐级放大至 0.98（fc）",
    ],
    fig="fig04_layer_sensitivity.png",
    fig_region=(6.55, 1.72, 6.5, 5.2),
    bullet_w=6.05)

# ---------------------------------------------------------------- slide 8
content_slide(
    "08", "四方法对比：Accuracy-α 全景",
    "Clean/Random-NAT/SGR-NAT 几乎重合；Fixed-NAT 全程平坦、唯一实质改善最坏情形。",
    [
        "α=+0.8：Clean 81.25% / Random 81.30% / SGR 82.06% / Fixed 91.79%",
        "SGR-NAT 与 Random-NAT 接近，与基线差异均在千分位量级",
        "Random/SGR 最佳验证轮次均为 epoch 0 → 训练信号被随机化稀释",
        "Fixed-NAT 曲线峰值移到 α=+0.4（94.27%）＝训练工作点对齐",
        "AURC：0.9283 / 0.9281 / 0.9290 / 0.9374",
    ],
    fig="fig03_alpha_sweep.png",
    fig_region=(6.45, 1.72, 6.6, 5.2),
    bullet_w=5.95)

# ---------------------------------------------------------------- slide 9
content_slide(
    "09", "Fixed-NAT 核心结果",
    "以 0.21 pp 干净精度代价，最差准确率 +10.54 pp、AURC 0.9283→0.9374。",
    [
        "Worst-Case Acc：81.25% → 91.79%（α=+0.8）",
        "AURC⁺（正半轴）：0.9128 → 0.9382；AURC⁻ 轻微让步 0.9399→0.9357",
        "Clean Accuracy 94.02%（基线 94.23%）",
        "方向不对称差 −0.0314 → +0.0007，正负风险对齐",
        "三个训练种子最差点 91.36%–91.79%（描述性一致，图10）",
    ],
    fig="fig06_worst_accuracy.png",
    fig_region=(6.45, 1.72, 6.6, 5.2),
    bullet_w=5.95)

# ---------------------------------------------------------------- slide 10
content_slide(
    "10", "校准分析：置信度崩溃型欠自信",
    "强正 α 使置信度崩溃（conf 0.25 / acc 81%）→ 欠自信而非过自信；Fixed-NAT 显著缓解。",
    [
        "α=+0.8（clean）：平均置信度 0.252、准确率 81.25%、ECE 0.560",
        "与 Guo 等报告的\u201c现代网络过自信\u201d方向相反：softmax 趋于均匀",
        "Fixed-NAT：置信度回升至 0.482、ECE 降至 0.436；训练点 α=+0.4 处 ECE 仅 0.028",
        "负 α 侧仅轻度过自信（conf≈0.99，ECE 0.05–0.07）",
        "安全含义：欠自信方向对拒识/降级决策相对保守，但仍需后校准",
    ],
    fig="fig09_reliability.png",
    fig_region=(6.3, 1.72, 6.75, 5.2),
    bullet_w=5.85)

# ---------------------------------------------------------------- slide 11
content_slide(
    "11", "消融：训练期 α 分布是决定因素",
    "把训练预算集中于最危险方向的单一工作点，优于摊薄到整个随机区间。",
    [
        "U(−0.5,+0.5) 全局随机（Random）：Worst 81.30%，α 期望为 0，训练信号被稀释",
        "灵敏度引导逐层 (0,α_g]（SGR）：Worst 82.06%，正向偏置带来轻微增益",
        "恒定 +0.4（Fixed）：Worst 91.79%，全部质量集中于危险方向中等强度",
        "best_epoch=0 现象：Random/SGR 微调几乎未能超过初始点",
        "适用前提：失效方向明确 + 部署期失真近似恒定（确定性失真 ≠ 零均值噪声）",
    ],
    fig="fig08_ece_alpha.png",
    fig_region=(6.55, 1.72, 6.5, 5.2),
    bullet_w=6.05)

# ---------------------------------------------------------------- slide 12
content_slide(
    "12", "结论与局限性",
    "深层激活非线性下，Fixed-NAT 是唯一实质有效的策略；结论限定于已实证的扰动口径。",
    [
        "结论 1：正 α 压缩为主导失效方向，根源是判别信息破坏而非扰动幅度",
        "结论 2：Fixed-NAT(+0.4) 最差准确率 +10.54 pp、AURC 0.9374、Clean 94.02%",
        "结论 3：强压缩引发欠自信崩溃，Fixed-NAT 将 ECE 0.560→0.436",
        "局限：注入范围为深层 4 层（4/21）；除 Fixed-NAT 外单种子、无显著性检验",
        "后续：修复命名碰撞后真全层注入复验、多种子重复、真实 CIM 硬件在环",
    ],
    fig="fig07_tradeoff.png",
    fig_region=(6.55, 1.72, 6.5, 5.2),
    bullet_w=6.05)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
